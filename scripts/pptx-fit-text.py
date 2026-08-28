#!/usr/bin/env python3
"""내보낸 PPTX 의 텍스트가 브라우저와 같은 자리에 머물게 한다.

브라우저와 파워포인트는 글자 폭을 미세하게 다르게 잰다. 이 저장소의 슬라이드는 예산을
꽉 채워 잡기 때문에, 그 미세한 차이가 곧 "한 줄 더"가 되고 상자를 넘친다. 두 가지를 건다.

1. **한 줄짜리 상자는 줄바꿈을 끈다.** HTML 의 제목은 대개 내용 크기로 잡히는 요소라
   애초에 감길 수 없다. 내보낸 상자는 그 폭을 그대로 쓰므로, 재는 방식이 조금만 달라도
   감기고 한 줄 높이 상자에서는 잘린다.

2. **여러 줄 상자에는 "넘치면 글자 축소"를 건다.** 파워포인트의 표준 동작(normAutofit)이라
   한 줄이 더 필요해지면 렌더러가 알아서 조금 줄여 카드 안에 담는다. 넘치지 않으면
   아무 일도 일어나지 않는다.

사용법: python3 scripts/pptx-fit-text.py <deck.pptx> [...]
"""
import glob
import os
import re
import sys
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import ImageFont

_CACHE = {}
# 서체가 ~/.fonts 에 없으면 폭을 잴 수 없어 이 스크립트가 아무것도 못 한다.
# 조용히 넘어가면 제목이 파워포인트에서 감기고 아무도 이유를 모른다 — 그래서 이름을 모아 알린다.
MISSING = set()


def _font(name, size):
    """설치된 실제 서체로 글자 폭을 잰다. 없으면 None — 그때는 축소를 걸지 않는다."""
    key = (name, round(size))
    if key in _CACHE:
        return _CACHE[key]
    want = (name or '').lower().replace(' ', '').replace('-', '')[:8]
    hit = None
    for f in glob.glob(os.path.expanduser('~/.fonts/*.ttf')):
        base = os.path.basename(f)[:-4].lower().replace(' ', '').replace('-', '')
        if want and base.startswith(want):
            hit = f
            break
    try:
        _CACHE[key] = ImageFont.truetype(hit, int(round(size))) if hit else None
    except Exception:
        _CACHE[key] = None
    if _CACHE[key] is None:
        MISSING.add(name or '(이름 없음)')
    return _CACHE[key]


def _line_height(tf, size):
    """줄간격은 추측하지 않고 문단에서 읽는다. 값이 없을 때만 1.35 로 본다."""
    for p_ in tf.paragraphs:
        if p_.line_spacing:
            return p_.line_spacing / 12700 if p_.line_spacing > 100 else p_.line_spacing * size
    return size * 1.35


def _needed_height(tf, width_pt, size):
    """상자 폭에서 실제로 몇 줄이 되는지 세어 필요한 높이를 낸다."""
    names = [r.font.name for p in tf.paragraphs for r in p.runs if r.font.name]
    fnt = _font(names[0] if names else '', size)
    if not fnt:
        return None
    line_h = _line_height(tf, size)
    lines = 0
    for para in tf.text.split('\n'):
        words = para.split(' ')
        cur, n = '', 1
        for w in words:
            t = (cur + ' ' + w).strip()
            if fnt.getlength(t) <= width_pt or not cur:
                cur = t
            else:
                n += 1
                cur = w
        lines += n
    return lines * line_h

def _collapse_whitespace(tf):
    """HTML 은 소스의 줄바꿈·들여쓰기를 공백 하나로 접지만, 내보낸 PPTX 는 그대로 싣는다.
    그래서 소스에서 들여쓴 줄이 슬라이드에서 들여쓴 채로 나온다. 같은 규칙으로 접는다."""
    n = 0
    for para in tf.paragraphs:
        runs = [r for r in para.runs]
        if not runs:
            continue
        for r in runs:
            new = re.sub(r'\s+', ' ', r.text)
            if new != r.text:
                r.text = new
                n += 1
        first, last = runs[0], runs[-1]
        if first.text.startswith(' '):
            first.text = first.text.lstrip(' ')
            n += 1
        if last.text.endswith(' '):
            last.text = last.text.rstrip(' ')
            n += 1
    return n


def _stays_on_slide(shape, tf, size, slide_w):
    """줄바꿈을 끄면 긴 줄이 슬라이드 밖으로 나갈 수 있다. 나갈 것 같으면 끄지 않는다."""
    names = [r.font.name for p in tf.paragraphs for r in p.runs if r.font.name]
    fnt = _font(names[0] if names else '', size)
    if not fnt:
        return False
    widest = max((fnt.getlength(line) for line in re.split(r'[\n\v]', tf.text)), default=0)
    # 렌더러가 우리보다 넓게 잴 수 있으므로 12% 여유를 두고 판단한다.
    return shape.left / 12700 + widest * 1.12 <= slide_w


def fix(path):
    prs = Presentation(path)
    slide_w = prs.slide_width / 12700
    nowrap = shrink = spaces = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text
            if not text.strip():
                continue
            spaces += _collapse_whitespace(tf)
            sizes = [r.font.size.pt for p in tf.paragraphs for r in p.runs if r.font.size]
            if not sizes:
                continue
            # 상자 높이가 담는 줄 수와 본문의 명시적 줄 수가 같으면, 그 줄바꿈은 작성자가
            # 정한 것이고 렌더러가 더 넣을 자리가 없다. 폭 계산이 브라우저보다 조금만 넓어도
            # 줄이 하나 늘어 상자를 넘치므로, 그런 상자는 자동 줄바꿈을 끈다.
            explicit = len(re.split(r'[\n\v]', text))
            line_h = _line_height(tf, max(sizes))
            fits = max(1, round(shape.height / 12700 / line_h))
            if explicit >= fits and _stays_on_slide(shape, tf, max(sizes), slide_w):
                tf.word_wrap = False
                nowrap += 1
            else:
                # 실제로 넘치는 상자에만 건다. 전부에 걸면 렌더러가 넘치지 않는 상자까지
                # 줄여서 카드마다 글자 크기가 들쭉날쭉해진다.
                need = _needed_height(tf, shape.width / 12700, max(sizes))
                # 딱 맞는 상자가 위험하다. 브라우저에서 정확히 들어차면 파워포인트에서는
                # 한 줄이 더 생겨 넘친다 — 92% 이상 찬 상자에 미리 축소를 걸어둔다.
                if need is not None and need > shape.height / 12700 * 0.92:
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    shrink += 1
    prs.save(path)
    return nowrap, shrink, spaces

if __name__ == '__main__':
    a = b = c = 0
    for p in sys.argv[1:]:
        x, y, z = fix(p)
        a += x; b += y; c += z
    print(f"줄바꿈 끔 {a}개 · 넘치면 축소 {b}개 · 공백 정리 {c}곳")
    if MISSING:
        print(f"  ⚠ 서체를 찾지 못해 폭을 재지 못했다: {', '.join(sorted(MISSING))}")
        print("    woff2 를 ttf 로 바꿔 ~/.fonts 에 넣고 fc-cache -f 를 돌린 뒤 다시 실행한다.")
